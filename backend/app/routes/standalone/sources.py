import os
import tempfile
import uuid
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form, Request
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select

from app.database import get_db
from app.services.auth_service import get_optional_user
from app.db_models import User, StandaloneSession, StandaloneSource
from app.models import StandaloneSourceResponse
from app.services import embedding_service
from app.services.website_service import fetch_webpage
from app.services.pdf_service import fetch_pdf
from app.utils.chunk_text import chunk_text

router = APIRouter()


def _get_guest_token(request: Request) -> str | None:
    return request.headers.get("X-Guest-Token")


async def _get_session_owner_check(
    db: AsyncSession, session_id: str, user: User | None, guest_token: str | None
) -> StandaloneSession:
    from app.routes.standalone.sessions import _get_session_owner_check as check
    return await check(db, session_id, user, guest_token)


def _make_index_key(session_id: str, source_id: str) -> str:
    return f"standalone_{session_id}_{source_id}"


def _source_to_response(src: StandaloneSource) -> StandaloneSourceResponse:
    return StandaloneSourceResponse(
        id=src.id,
        session_id=src.session_id,
        source_type=src.source_type,
        title=src.title,
        metadata_json=src.metadata_json,
        file_name=src.file_name,
        created_at=src.created_at.isoformat() if src.created_at else "",
    )


@router.post("/{session_id}/sources", response_model=StandaloneSourceResponse, status_code=201)
async def upload_source_text(
    session_id: str,
    request: Request,
    source_type: str = Form("text"),
    title: str = Form(""),
    content: str = Form(""),
    url: str = Form(""),
    file: UploadFile | None = None,
    user: User | None = Depends(get_optional_user),
    db: AsyncSession = Depends(get_db),
):
    guest_token = _get_guest_token(request)
    await _get_session_owner_check(db, session_id, user, guest_token)

    actual_type = source_type
    actual_title = title
    actual_content = content
    file_name = None

    if source_type == "url":
        if not url:
            raise HTTPException(status_code=422, detail={"error": "URL_REQUIRED", "message": "URL is required for url type."})
        try:
            page_data = await fetch_webpage(url)
            actual_content = page_data.text
            actual_title = page_data.title or title or url
        except Exception as e:
            raise HTTPException(status_code=422, detail={"error": "FETCH_FAILED", "message": f"Failed to fetch URL: {str(e)}"})
        actual_type = "website"

    elif source_type == "file":
        if not file:
            raise HTTPException(status_code=422, detail={"error": "FILE_REQUIRED", "message": "File is required for file type."})
        file_name = file.filename or "unknown"
        raw = await file.read()
        ext = os.path.splitext(file_name)[1].lower()

        if ext in (".pdf",):
            import fitz
            doc = fitz.open(stream=raw, filetype="pdf")
            pages = []
            for page in doc:
                pages.append(page.get_text())
            actual_content = "\n\n".join(pages)
            actual_title = title or file_name
        elif ext in (".docx",):
            import docx
            from io import BytesIO
            doc = docx.Document(BytesIO(raw))
            actual_content = "\n".join(p.text for p in doc.paragraphs)
            actual_title = title or file_name
        elif ext in (".pptx",):
            from pptx import Presentation
            from io import BytesIO
            prs = Presentation(BytesIO(raw))
            slides = []
            for slide in prs.slides:
                slide_text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                slides.append(slide_text)
            actual_content = "\n\n".join(slides)
            actual_title = title or file_name
        else:
            actual_content = raw.decode("utf-8", errors="replace")
            actual_title = title or file_name

        actual_type = ext.lstrip(".") if ext else "text"

    elif source_type == "text":
        if not actual_content.strip():
            raise HTTPException(status_code=422, detail={"error": "CONTENT_REQUIRED", "message": "Content is required for text type."})
        actual_title = title or "Untitled"

    if not actual_content.strip():
        raise HTTPException(status_code=422, detail={"error": "NO_CONTENT", "message": "No content could be extracted."})

    source_id = str(uuid.uuid4())
    index_key = _make_index_key(session_id, source_id)

    source = StandaloneSource(
        id=source_id,
        session_id=session_id,
        source_type=actual_type,
        title=actual_title,
        content=actual_content,
        index_key=index_key,
        file_name=file_name,
    )
    db.add(source)
    await db.flush()

    chunks = chunk_text(actual_content, 500, 50)
    await embedding_service.index_transcript(index_key, actual_content)

    await db.commit()
    await db.refresh(source)

    return _source_to_response(source)


@router.get("/{session_id}/sources", response_model=list[StandaloneSourceResponse])
async def list_sources(
    session_id: str,
    request: Request,
    user: User | None = Depends(get_optional_user),
    db: AsyncSession = Depends(get_db),
):
    guest_token = _get_guest_token(request)
    await _get_session_owner_check(db, session_id, user, guest_token)

    result = await db.execute(
        select(StandaloneSource).where(StandaloneSource.session_id == session_id)
    )
    sources = result.scalars().all()
    return [_source_to_response(s) for s in sources]


@router.delete("/{session_id}/sources/{source_id}")
async def delete_source(
    session_id: str,
    source_id: str,
    request: Request,
    user: User | None = Depends(get_optional_user),
    db: AsyncSession = Depends(get_db),
):
    guest_token = _get_guest_token(request)
    await _get_session_owner_check(db, session_id, user, guest_token)

    result = await db.execute(
        select(StandaloneSource).where(
            StandaloneSource.id == source_id,
            StandaloneSource.session_id == session_id,
        )
    )
    src = result.scalar_one_or_none()
    if not src:
        raise HTTPException(status_code=404, detail={"error": "NOT_FOUND", "message": "Source not found."})

    embedding_service.delete_chunks(src.index_key)
    await db.delete(src)
    await db.commit()

    return {"status": "deleted"}

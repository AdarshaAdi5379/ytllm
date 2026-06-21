import hashlib
import io

from pptx import Presentation


PPTX_MAX_SIZE = 50 * 1024 * 1024


class PptxResult:
    def __init__(self, title: str, text: str, index_key: str):
        self.title = title
        self.text = text
        self.index_key = index_key


def content_to_index_key(content: bytes) -> str:
    return "pptx_" + hashlib.sha256(content).hexdigest()[:16]


def process_pptx(file_bytes: bytes, title: str = "") -> PptxResult:
    if len(file_bytes) > PPTX_MAX_SIZE:
        raise ValueError(f"PPTX exceeds size limit of {PPTX_MAX_SIZE // (1024 * 1024)}MB.")

    prs = Presentation(io.BytesIO(file_bytes))

    blocks: list[str] = []
    first_title = ""

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_blocks: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                slide_blocks.append(text)

        if slide_blocks:
            slide_text = "\n".join(slide_blocks)
            blocks.append(f"[Slide {slide_num}]\n{slide_text}")
            if slide_num == 1 and not first_title:
                first_title = slide_blocks[0][:200]

    if hasattr(prs, 'notes_master') and prs.notes_master:
        for slide_num, slide in enumerate(prs.slides, 1):
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    blocks.append(f"[Slide {slide_num} Notes]\n{notes}")

    full_text = "\n\n".join(blocks)
    if len(full_text) < 50:
        raise ValueError("Could not extract meaningful content from this PPTX file.")

    title = title.strip() or first_title or "Untitled Presentation"
    index_key = content_to_index_key(file_bytes)

    return PptxResult(title=title, text=full_text, index_key=index_key)
